import docx2txt
import os
import tempfile
from django.contrib import messages
from django.contrib.auth.decorators import login_required
from django.http import JsonResponse, HttpResponseNotFound
from django.shortcuts import render, redirect, get_object_or_404
from dotenv import load_dotenv

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings
from langchain.vectorstores.faiss import FAISS
from langchain.chains import ConversationalRetrievalChain
from langchain.callbacks import get_openai_callback
from langchain.chat_models import ChatOpenAI
from langchain.memory import ConversationBufferMemory


from PyPDF2 import PdfFileReader
from pptx import Presentation
from .forms import PDFUploadForm, PDFUpdateForm, PDFDocumentForm2
from .models import ChatMessage, PDFDocument, UserData

load_dotenv()


def main(request):
    """
    The main function is the entry point for the chat_llm app.
    
    :param request: Get information about the user who is currently logged in
    :return: A render of the index
    """
    # avatar = Avatar.objects.filter(user_id=request.user.id).first()
    return render(request, 'chat_llm/index.html', context={})


def get_pdf_text(file):
    """
    Retrieves text from a file.

    :param file: File object.
    :return: Extracted text from the file.
    """
    text = None
    if file:
        with tempfile.NamedTemporaryFile(delete=False) as temp_file:
            temp_file.write(file.read())
            temp_file.flush()
            if file.name.endswith('.pdf'):
                pdf_reader = PdfFileReader(temp_file.name)
                text = ''.join(page.extract_text() for page in pdf_reader.pages)
            elif file.name.endswith('.txt'):
                with open(temp_file.name, 'r') as f:
                    text = f.read()
            elif file.name.endswith('.docx'):
                text = docx2txt.process(temp_file.name)
            elif file.name.endswith('.pptx'):
                prs = Presentation(temp_file.name)
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                text = '\n'.join(text_runs)
    return text


def get_text_chunks(text):
    """
    Splits text into chunks.

    :param text: Input text.
    :return: List of text chunks.
    """

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    chunks = text_splitter.split_text(text)
    return chunks


def get_vectorstore(text_chunks):
    """
    Retrieves the vector store for text chunks.

    :param text_chunks: List of text chunks.
    :return: Knowledge base vector store.
    """
    api_key = os.getenv("OPENAI_API_KEY")
    embeddings = OpenAIEmbeddings(model="text-embedding-ada-002", openai_api_key=api_key)

    vectorstore = FAISS.from_texts(text_chunks, embeddings)
    return vectorstore


def get_conversation_chain(vectorstore):
    """
    Retrieves the conversation chain.

    :param vectorstore: Vector store.
    :return: Conversational retrieval chain.
    """
    llm = ChatOpenAI(model="gpt-3.5-turbo")
    memory = ConversationBufferMemory(memory_key='chat_history', return_messages=True)
    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory
    )
    return conversation_chain


@login_required(login_url="/login/")
def upload_pdf(request):
    """
    Handles the PDF upload.

    :param request: HTTP request.
    :return: JSON response.
    """
    user = request.user

    try:
        user_data = UserData.objects.get(user=user)
    except UserData.DoesNotExist:
        # Обробка ситуації, коли відсутній запис UserData для користувача.
        # Можна створити запис за замовчуванням тут.
        user_data = UserData.objects.create(user=user, subscribe_plan='free', total_files_uploaded=0,
                                            total_questions_asked=0)

    # Отримати обмеження для користувача залежно від плану підписки
    max_files_allowed = user_data.max_files_allowed_for_plan()

    if request.method == 'POST':
        form = PDFUploadForm(request.POST, request.FILES)
        if form.is_valid():
            pdf_document = request.FILES['pdf_document']


            # Перевірка розширення файлу
            _, file_extension = os.path.splitext(pdf_document.name)
            if file_extension.lower() not in [".pdf", ".txt", ".docx", ".pptx"]:
                return JsonResponse({'error': 'Only PDF, TXT, DOCX, PPTX files'}, status=400)

            if pdf_document.size > 50 * 1024 * 1024:  # Розмір файлу понад 50 МБ
                return JsonResponse({'error': "File size exceeds 50 MB."}, status=400)

            # Check if the user has exceeded the limit for uploaded files
            if user_data.total_files_uploaded >= max_files_allowed:
                return JsonResponse({'error': 'You have reached the limit for uploaded files.'}, status=400)

            # Check if a file with the same name already exists for this user
            if PDFDocument.objects.filter(user=user, title=pdf_document.name).exists():
                return JsonResponse({'error': 'A file with the same name already exists.'}, status=400)

            # Збільшити кількість завантажених файлів користувача
            user_data.total_files_uploaded += 1
            user_data.save()

            # Зберегти PDF-документ у базі даних
            pdf = PDFDocument(user=user, title=pdf_document.name)
            pdf.documentContent = get_pdf_text(pdf_document)
            pdf.save()
            return JsonResponse({'success': 'File uploaded successfully.'})

    else:
        form = PDFUploadForm()
    user_pdfs = PDFDocument.objects.filter(user=request.user)
    chat_message = ChatMessage.objects.all()
    return render(request, 'chat_llm/chat_base.html', {'form': form, 'user_pdfs': user_pdfs})


@login_required(login_url="/login/")
def ask_question(request):
    """
    Handles the user's question and generates a response.

    :param request: HTTP request.
    :return: Rendered page with the response.
    """
    try:
        user = request.user

        try:
            user_data = UserData.objects.get(user=user)
        except UserData.DoesNotExist:
            # Обробка ситуації, коли відсутній запис UserData для користувача.
            # Можна створити запис за замовчуванням тут.
            user_data = UserData.objects.create(user=user, subscribe_plan='free', total_files_uploaded=0,
                                                total_questions_asked=0)

        max_questions_allowed = user_data.max_questions_allowed_for_plan()

        chat_history = ChatMessage.objects.filter(user=request.user).order_by(
            'timestamp')[:10]
        chat_response = ''
        user_pdfs = PDFDocument.objects.filter(user=request.user)
        user_question = ""
        selected_pdf = None  # Змінено з selected_pdf_id на об'єкт PDFDocument

        if request.method == 'POST':
            # Check if the user has exceeded the limit for questions per file
            if user_data.total_questions_asked >= max_questions_allowed:
                return JsonResponse({'error': 'You have reached the limit for questions.'}, status=400)
            else:
                user_question = request.POST.get('user_question')
                selected_pdf_id = request.POST.get('selected_pdf')
                selected_pdf = get_object_or_404(PDFDocument, id=selected_pdf_id)
                text_chunks = get_text_chunks(selected_pdf.documentContent)

                knowledge_base = get_vectorstore(text_chunks)
                conversation_chain = get_conversation_chain(knowledge_base)

                with get_openai_callback() as cb:
                    response = conversation_chain({'question': user_question})

                chat_response = response["answer"]
                chat_message = ChatMessage(user=request.user, message=user_question, answer=chat_response,
                                           pdf_document=selected_pdf)  # Передаємо об'єкт PDFDocument
                user_data.total_questions_asked += 1
                user_data.save()
                chat_message.save()
    except MemoryError as e:
        # Обробити MemoryError тут
        error_message = ("Недостатньо пам'яті на безкоштовному сервісі Fjy.io для обробки запиту. "
                         "Спробуйте розгорнути локально з скрипту на GitHub.")
        context = {'error_message': error_message}
        return render(request, 'chat_llm/chat_base.html', context)
    except Exception as e:
        # Обробка інших винятків
        error_message = f"Сталася помилка: {str(e)}"
        context = {'error_message': error_message}
        return render(request, 'chat_llm/chat_base.html', context)

    # Отримуємо повідомлення, які відносяться до обраного PDFDocument
    chat_message = ChatMessage.objects.filter(user=request.user, pdf_document=selected_pdf).order_by('timestamp')

    context = {'chat_response': chat_response, 'chat_history': chat_history, 'user_question': user_question,
               'user_pdfs': user_pdfs, 'chat_message': chat_message}

    return render(request, 'chat_llm/chat_base.html', context)


@login_required(login_url="/login/")
def get_chat_history(request):
    """
    The get_chat_history function returns a JSON response containing the chat history for a given PDF document.
    
    :param request: Get the request object
    :return: A list of dictionaries containing the chat history for a given PDF document
    """
    pdf_id = request.GET.get('pdf_id')
    if not pdf_id:
        return JsonResponse({"error": "PDF ID not provided"}, status=400)

    # Тут ви отримуєте історію чату для даного pdf_id
    # Це припущення, що у вас є модель ChatMessage або щось подібне
    messages = ChatMessage.objects.filter(pdf_document_id=pdf_id).values('message', 'answer', 'timestamp')

    return JsonResponse(list(messages), safe=False)


@login_required(login_url="/login/")
def view_pdf(request, pdf_id):
    """
    Displays a PDF document.

    :param request: HTTP request.
    :param pdf_id: ID of the PDF document.
    :return: Rendered page with the PDF document.
    """
    pdf = PDFDocument.objects.get(id=pdf_id)
    return render(request, 'view_pdf.html', {'pdf': pdf})


@login_required(login_url="/login/")
def view_chat_history(request):
    """
    Displays the chat history for the logged-in user.

    :param request: HTTP request.
    :return: Rendered page with chat history.
    """
    chat_messages = ChatMessage.objects.filter(user=request.user)
    return render(request, 'view_chat_history.html', {'chat_messages': chat_messages})


def list_pdfs(request):
    """
    Lists PDF documents for the logged-in user.

    :param request: HTTP request.
    :return: Rendered page with the list of PDF documents.
    """
    pdfs = PDFDocument.objects.filter(user=request.user)
    return render(request, 'edit_pdf.html', {'pdfs': pdfs})


def delete_pdf(request, pdf_id):
    """
    Deletes a PDF document.

    :param request: HTTP request.
    :param pdf_id: ID of the PDF document to delete.
    :return: Redirect or error response.
    """
    try:
        pdfs = PDFDocument.objects.get(id=pdf_id)
        pdf_title = getattr(pdfs, 'document', pdfs.title)
        pdfs.delete()
        return redirect('/pdfs/')
    except PDFDocument.DoesNotExist:
        # Handle the case where the PDFDocument with the given id does not exist
        return HttpResponseNotFound("PDF not found")


def update_pdf(request, pdf_id):
    """
    Updates a PDF document.

    :param request: HTTP request.
    :param pdf_id: ID of the PDF document to update.
    :return: Rendered page with the updated PDF document.
    """
    pdf = get_object_or_404(PDFDocument, pk=pdf_id)
    if request.method == 'POST':
        form = PDFDocumentForm2(request.POST, instance=pdf)
        if form.is_valid():
            form.save()
            messages.success(request, 'Pdf updated successfully.')
            return redirect('/pdfs/')
    else:
        form = PDFUpdateForm(instance=pdf)
    return render(request, 'update_pdf.html', {'form': form, 'pdf': pdf})
